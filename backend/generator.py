import os
import json
import tempfile
import unicodedata
from datetime import datetime
from dotenv import load_dotenv
from groq import Groq
from fastapi.responses import FileResponse
from fastapi import BackgroundTasks

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls

from fpdf import FPDF

load_dotenv()

client = Groq(api_key=os.getenv("GROQ_API_KEY"))

PRIMARY_MODEL = os.getenv("GROQ_MODEL", "openai/gpt-oss-120b")
FALLBACK_MODEL = "openai/gpt-oss-20b"

# ─── Color Themes Per Document Type ───────────────────────────────────────────

# Themes for Presentations (Full slide backgrounds & cards)
PPT_THEMES = {
    "dark": {
        "bg": (14, 14, 16),
        "card_bg": (24, 24, 28),
        "border": (45, 45, 52),
        "accent": (232, 160, 32),       # Amber
        "title": (255, 255, 255),
        "text": (190, 195, 205),
        "muted": (130, 135, 145),
    },
    "blue": {
        "bg": (10, 18, 34),
        "card_bg": (16, 28, 52),
        "border": (30, 50, 85),
        "accent": (59, 130, 246),       # Modern Blue
        "title": (255, 255, 255),
        "text": (205, 220, 245),
        "muted": (130, 155, 190),
    },
    "green": {
        "bg": (10, 24, 16),
        "card_bg": (18, 38, 26),
        "border": (28, 62, 42),
        "accent": (20, 184, 166),       # Teal
        "title": (255, 255, 255),
        "text": (200, 235, 220),
        "muted": (125, 175, 155),
    },
    "purple": {
        "bg": (20, 14, 32),
        "card_bg": (32, 22, 50),
        "border": (54, 38, 82),
        "accent": (168, 85, 247),       # Purple
        "title": (255, 255, 255),
        "text": (230, 215, 250),
        "muted": (155, 135, 185),
    },
    "light": {
        "bg": (248, 249, 250),
        "card_bg": (255, 255, 255),
        "border": (220, 225, 232),
        "accent": (217, 119, 6),        # Warm Amber for light mode
        "title": (20, 24, 32),
        "text": (70, 80, 95),
        "muted": (120, 130, 145),
    },
}

# Themes for Word Documents (Professional white/light paper with theme accents)
DOC_THEMES = {
    "dark": {
        "accent": DocxRGB(217, 119, 6),      # Deep Amber Accent
        "title": DocxRGB(17, 24, 39),        # Charcoal
        "heading": DocxRGB(217, 119, 6),    # Amber
        "text": DocxRGB(55, 65, 81),         # High-contrast readable body text
        "muted": DocxRGB(107, 114, 128),     # Slate Grey
        "box_bg": "F3F4F6",                  # Clean Light Gray Callout
        "box_border": "D97706",
    },
    "blue": {
        "accent": DocxRGB(37, 99, 235),      # Royal Blue
        "title": DocxRGB(15, 23, 42),        # Dark Navy
        "heading": DocxRGB(37, 99, 235),     # Royal Blue
        "text": DocxRGB(51, 65, 85),         # Slate body text
        "muted": DocxRGB(100, 116, 139),
        "box_bg": "EFF6FF",                  # Light Blue Callout
        "box_border": "2563EB",
    },
    "green": {
        "accent": DocxRGB(13, 148, 136),     # Teal
        "title": DocxRGB(6, 78, 59),         # Deep Emerald
        "heading": DocxRGB(13, 148, 136),    # Teal
        "text": DocxRGB(55, 65, 81),
        "muted": DocxRGB(100, 116, 139),
        "box_bg": "F0FDF4",                  # Light Emerald Callout
        "box_border": "0D9488",
    },
    "purple": {
        "accent": DocxRGB(124, 58, 237),     # Violet
        "title": DocxRGB(59, 7, 100),        # Deep Purple
        "heading": DocxRGB(124, 58, 237),    # Violet
        "text": DocxRGB(55, 65, 81),
        "muted": DocxRGB(107, 114, 128),
        "box_bg": "FAF5FF",                  # Light Violet Callout
        "box_border": "7C3AED",
    },
    "light": {
        "accent": DocxRGB(71, 85, 105),      # Slate
        "title": DocxRGB(15, 23, 42),
        "heading": DocxRGB(51, 65, 85),
        "text": DocxRGB(55, 65, 81),
        "muted": DocxRGB(148, 163, 184),
        "box_bg": "F8FAFC",
        "box_border": "64748B",
    },
}

# Themes for PDF Documents (Consistent background on every page)
PDF_THEMES = {
    "dark": {
        "is_dark": True,
        "bg": (14, 14, 18),
        "card_bg": (22, 22, 28),
        "border": (45, 45, 56),
        "accent": (232, 160, 32),       # Amber
        "title": (255, 255, 255),
        "heading": (232, 160, 32),
        "text": (215, 220, 230),        # Crisp High-Contrast Light Grey
        "muted": (140, 145, 160),
    },
    "blue": {
        "is_dark": True,
        "bg": (10, 18, 34),
        "card_bg": (16, 28, 52),
        "border": (30, 52, 90),
        "accent": (59, 130, 246),
        "title": (255, 255, 255),
        "heading": (96, 165, 250),
        "text": (215, 230, 250),
        "muted": (140, 165, 195),
    },
    "green": {
        "is_dark": True,
        "bg": (10, 24, 18),
        "card_bg": (18, 38, 28),
        "border": (28, 65, 45),
        "accent": (20, 184, 166),
        "title": (255, 255, 255),
        "heading": (45, 212, 191),
        "text": (210, 240, 230),
        "muted": (135, 175, 160),
    },
    "purple": {
        "is_dark": True,
        "bg": (20, 14, 32),
        "card_bg": (32, 22, 50),
        "border": (56, 40, 84),
        "accent": (168, 85, 247),
        "title": (255, 255, 255),
        "heading": (192, 132, 252),
        "text": (235, 225, 250),
        "muted": (160, 140, 190),
    },
    "light": {
        "is_dark": False,
        "bg": (255, 255, 255),
        "card_bg": (248, 250, 252),
        "border": (226, 232, 240),
        "accent": (217, 119, 6),
        "title": (15, 23, 42),
        "heading": (217, 119, 6),
        "text": (51, 65, 85),
        "muted": (100, 116, 139),
    },
}


# ─── Robust Unicode Sanitizer for PDF ─────────────────────────────────────────
def clean_pdf_text(text: str) -> str:
    """Normalizes all unicode characters (dashes, quotes, spaces, math symbols) to safe ASCII."""
    if not text or not isinstance(text, str):
        return ""
    
    char_map = {
        "\u2018": "'", "\u2019": "'", "\u201a": "'", "\u201b": "'",
        "\u201c": '"', "\u201d": '"', "\u201e": '"', "\u201f": '"',
        "\u2013": "-", "\u2014": " - ", "\u2015": " - ", "\u2010": "-", "\u2011": "-", "\u2012": "-", "\u2212": "-",
        "\u2026": "...", "\u2022": "*", "\u00b7": "*", "\u25cf": "*", "\u25cb": "*",
        "\u00d7": "x", "\u2715": "x", "\u2716": "x",
        "\u00a0": " ", "\u202f": " ", "\u2009": " ", "\u200a": " ", "\u200b": "", "\ufeff": "",
        "\u2192": "->", "\u2190": "<-", "\u21d2": "=>", "\u203a": ">", "\u00bb": ">>", "\u00ab": "<<",
        "\u2002": " ", "\u2003": " ", "\u2004": " ", "\u2005": " ", "\u2006": " ", "\u2007": " ", "\u2008": " ",
        "\u2044": "/", "\u2215": "/", "\u2032": "'", "\u2033": '"',
        "\u2264": "<=", "\u2265": ">=", "\u2260": "!=", "\u2248": "~",
    }
    for orig, repl in char_map.items():
        text = text.replace(orig, repl)
        
    text = unicodedata.normalize("NFKD", text)
    return text.encode("ascii", "ignore").decode("ascii")


# ─── Robust AI Query Helper ───────────────────────────────────────────────────
def query_groq(prompt: str) -> dict:
    try:
        response = client.chat.completions.create(
            model=PRIMARY_MODEL,
            response_format={"type": "json_object"},
            messages=[{"role": "user", "content": prompt}],
            temperature=0.6,
        )
    except Exception as e:
        print(f"Primary model {PRIMARY_MODEL} failed ({e}), using {FALLBACK_MODEL}")
        response = client.chat.completions.create(
            model=FALLBACK_MODEL,
            response_format={"type": "json_object"},
            messages=[{"role": "user", "content": prompt}],
            temperature=0.6,
        )

    raw = response.choices[0].message.content.strip()
    return json.loads(raw)


# ─── 1. POWERPOINT GENERATION (Structured 16:9 Multi-Card Deck) ───────────────
def get_ppt_ai_content(topic: str, slide_count: int) -> dict:
    prompt = f"""
You are an executive presentation designer. Create an exhaustive, highly structured presentation deck on: "{topic}".
Target length: {slide_count} slides total.

Return ONLY a JSON object with this exact structure:
{{
  "deck_title": "Concise Main Title",
  "deck_subtitle": "Comprehensive professional subtitle explaining the core objective",
  "category": "CATEGORY TAG (e.g. SYSTEM ARCHITECTURE / STRATEGIC OVERVIEW)",
  "slides": [
    {{
      "title": "Specific Slide Title",
      "category": "SECTION TAG",
      "subtitle": "Informative single-sentence subheader",
      "cards": [
        {{
          "heading": "Core Sub-topic 1",
          "detail": "2-3 sentences of rich, practical, in-depth explanation with key technical or strategic insights."
        }},
        {{
          "heading": "Core Sub-topic 2",
          "detail": "2-3 sentences of rich, practical, in-depth explanation with key technical or strategic insights."
        }},
        {{
          "heading": "Core Sub-topic 3",
          "detail": "2-3 sentences of rich, practical, in-depth explanation with key technical or strategic insights."
        }}
      ],
      "takeaway": "Key Takeaway: Clear, actionable summary of this slide's core message."
    }}
  ]
}}
Generate exactly {slide_count - 1} content slides (the title slide is created automatically from deck_title/subtitle).
Ensure details are rich, informative, and professional. No placeholder text.
"""
    return query_groq(prompt)


async def generate_ppt(req, background_tasks: BackgroundTasks):
    data = get_ppt_ai_content(req.topic, req.slide_count)
    theme = PPT_THEMES.get(req.theme, PPT_THEMES["dark"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Executive Title Slide ────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = RGBColor(*theme["bg"])

    # Top Category Pill Box
    cat_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.5))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = data.get("category", "EXECUTIVE BRIEFING").upper()
    p_cat.runs[0].font.size = Pt(11)
    p_cat.runs[0].font.bold = True
    p_cat.runs[0].font.color.rgb = RGBColor(*theme["accent"])

    # Main Title
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = data.get("deck_title", req.topic)
    p_title.runs[0].font.size = Pt(40)
    p_title.runs[0].font.bold = True
    p_title.runs[0].font.color.rgb = RGBColor(*theme["title"])

    # Subtitle
    sub_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = data.get("deck_subtitle", "Comprehensive Analysis and Practical Implementation Framework")
    p_sub.runs[0].font.size = Pt(18)
    p_sub.runs[0].font.color.rgb = RGBColor(*theme["text"])

    # Title Slide Footer Meta
    meta_box = s1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.6))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = f"DocuCraft Document Engine  •  {datetime.now().strftime('%B %Y')}  •  {req.slide_count} Slides"
    p_meta.runs[0].font.size = Pt(11)
    p_meta.runs[0].font.color.rgb = RGBColor(*theme["muted"])

    # ── Content Slides: 3-Card Grid with Takeaway Banner ──────────────────────
    content_slides = data.get("slides", [])
    total_slides = len(content_slides) + 1

    for idx, s_data in enumerate(content_slides, start=2):
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(*theme["bg"])

        # Top Category & Slide Title
        header_box = s.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.2))
        tf_h = header_box.text_frame
        tf_h.word_wrap = True

        p_hcat = tf_h.paragraphs[0]
        p_hcat.text = s_data.get("category", "ANALYSIS").upper()
        p_hcat.runs[0].font.size = Pt(10)
        p_hcat.runs[0].font.bold = True
        p_hcat.runs[0].font.color.rgb = RGBColor(*theme["accent"])

        p_htitle = tf_h.add_paragraph()
        p_htitle.text = s_data.get("title", f"Slide {idx}")
        p_htitle.runs[0].font.size = Pt(22)
        p_htitle.runs[0].font.bold = True
        p_htitle.runs[0].font.color.rgb = RGBColor(*theme["title"])

        if s_data.get("subtitle"):
            p_hsub = tf_h.add_paragraph()
            p_hsub.text = s_data.get("subtitle")
            p_hsub.runs[0].font.size = Pt(12)
            p_hsub.runs[0].font.color.rgb = RGBColor(*theme["text"])

        # 3 Cards Grid
        cards = s_data.get("cards", [])[:3]
        card_w = Inches(3.68)
        card_h = Inches(3.5)
        card_gap = Inches(0.33)
        start_x = Inches(0.8)
        card_y = Inches(1.95)

        for c_idx, card in enumerate(cards):
            cx = start_x + (c_idx * (card_w + card_gap))
            
            # Card Background Shape
            card_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_w, card_h)
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(*theme["card_bg"])
            card_shape.line.color.rgb = RGBColor(*theme["border"])
            card_shape.line.width = Pt(1)

            # Card Text
            tf_card = card_shape.text_frame
            tf_card.vertical_anchor = MSO_ANCHOR.TOP
            tf_card.word_wrap = True
            tf_card.margin_left = Inches(0.25)
            tf_card.margin_right = Inches(0.25)
            tf_card.margin_top = Inches(0.25)
            tf_card.margin_bottom = Inches(0.25)

            p_chead = tf_card.paragraphs[0]
            p_chead.text = card.get("heading", f"Point {c_idx + 1}")
            p_chead.runs[0].font.size = Pt(14)
            p_chead.runs[0].font.bold = True
            p_chead.runs[0].font.color.rgb = RGBColor(*theme["accent"])
            p_chead.space_after = Pt(8)

            p_cdetail = tf_card.add_paragraph()
            p_cdetail.text = card.get("detail", "")
            p_cdetail.runs[0].font.size = Pt(11)
            p_cdetail.runs[0].font.color.rgb = RGBColor(*theme["text"])
            p_cdetail.line_spacing = 1.25

        # Bottom Takeaway Box
        takeaway_text = s_data.get("takeaway", "")
        if takeaway_text:
            take_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.75), Inches(11.73), Inches(0.85))
            take_shape.fill.solid()
            take_shape.fill.fore_color.rgb = RGBColor(*theme["card_bg"])
            take_shape.line.color.rgb = RGBColor(*theme["accent"])
            take_shape.line.width = Pt(1)

            tf_take = take_shape.text_frame
            tf_take.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf_take.margin_left = Inches(0.3)
            p_take = tf_take.paragraphs[0]
            p_take.text = takeaway_text
            p_take.runs[0].font.size = Pt(11)
            p_take.runs[0].font.color.rgb = RGBColor(*theme["title"])

        # Slide Number Footer
        foot_box = s.shapes.add_textbox(Inches(10.0), Inches(6.8), Inches(2.5), Inches(0.4))
        p_foot = foot_box.text_frame.paragraphs[0]
        p_foot.alignment = PP_ALIGN.RIGHT
        p_foot.text = f"{idx:02d} / {total_slides:02d}"
        p_foot.runs[0].font.size = Pt(10)
        p_foot.runs[0].font.color.rgb = RGBColor(*theme["muted"])

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    background_tasks.add_task(os.unlink, tmp.name)
    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{req.topic}.pptx",
        background=background_tasks,
    )


# ─── 2. WORD DOCUMENT GENERATION (Technical Whitepaper) ───────────────────────
def get_doc_ai_content(topic: str, section_count: int) -> dict:
    prompt = f"""
You are a senior technical writer and research analyst. Produce an in-depth, publication-ready technical whitepaper on: "{topic}".
Target: {section_count} full sections.

Return ONLY a JSON object with this exact structure:
{{
  "title": "Comprehensive Document Title",
  "subtitle": "Executive subtitle detailing the scope and significance",
  "executive_summary": "A rich 3-4 sentence executive overview explaining the subject, relevance, and core findings.",
  "sections": [
    {{
      "heading": "1. Section Title",
      "intro": "A thorough, 4-5 sentence analytical paragraph breaking down the core concepts, mechanisms, and background context in depth.",
      "key_points": [
        {{
          "title": "Point Concept",
          "description": "2-3 sentences explaining the mechanism, impact, and practical considerations."
        }},
        {{
          "title": "Point Concept",
          "description": "2-3 sentences explaining the mechanism, impact, and practical considerations."
        }},
        {{
          "title": "Point Concept",
          "description": "2-3 sentences explaining the mechanism, impact, and practical considerations."
        }}
      ],
      "takeaway": "Key Recommendation: Concise strategic insight for engineering or business teams."
    }}
  ],
  "conclusion": "A comprehensive concluding summary covering future outlook and best practices."
}}
Generate exactly {section_count} sections. Every section must have deep, substantial, educational content without filler.
"""
    return query_groq(prompt)


def set_cell_background(cell, fill_hex: str):
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.append(parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>'))


async def generate_doc(req, background_tasks: BackgroundTasks):
    data = get_doc_ai_content(req.topic, req.slide_count)
    theme = DOC_THEMES.get(req.theme, DOC_THEMES["dark"])

    doc = Document()

    # Set 1-inch margins
    for section in doc.sections:
        section.top_margin = DocxInches(1.0)
        section.bottom_margin = DocxInches(1.0)
        section.left_margin = DocxInches(1.0)
        section.right_margin = DocxInches(1.0)

    # ── Title & Subtitle ──────────────────────────────────────────────────────
    p_title = doc.add_paragraph()
    p_title.paragraph_format.space_before = DocxPt(0)
    p_title.paragraph_format.space_after = DocxPt(4)
    r_title = p_title.add_run(data.get("title", req.topic))
    r_title.font.name = "Arial"
    r_title.font.size = DocxPt(24)
    r_title.font.bold = True
    r_title.font.color.rgb = theme["title"]

    p_sub = doc.add_paragraph()
    p_sub.paragraph_format.space_after = DocxPt(14)
    r_sub = p_sub.add_run(data.get("subtitle", "Technical Overview and Strategic Insights"))
    r_sub.font.name = "Arial"
    r_sub.font.size = DocxPt(12)
    r_sub.font.color.rgb = theme["muted"]

    # ── Metadata Line ─────────────────────────────────────────────────────────
    p_meta = doc.add_paragraph()
    p_meta.paragraph_format.space_after = DocxPt(16)
    r_meta = p_meta.add_run(f"Prepared by DocuCraft AI Engine  |  {datetime.now().strftime('%B %d, %Y')}  |  Confidential")
    r_meta.font.size = DocxPt(9.5)
    r_meta.font.italic = True
    r_meta.font.color.rgb = theme["muted"]

    # ── Executive Summary Callout Box ─────────────────────────────────────────
    exec_summary = data.get("executive_summary", "")
    if exec_summary:
        tbl = doc.add_table(rows=1, cols=1)
        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = tbl.cell(0, 0)
        set_cell_background(cell, theme["box_bg"])
        
        cp = cell.paragraphs[0]
        cp.paragraph_format.space_before = DocxPt(8)
        cp.paragraph_format.space_after = DocxPt(6)
        r_ex_lbl = cp.add_run("EXECUTIVE SUMMARY\n")
        r_ex_lbl.font.size = DocxPt(10)
        r_ex_lbl.font.bold = True
        r_ex_lbl.font.color.rgb = theme["accent"]

        r_ex_txt = cp.add_run(exec_summary)
        r_ex_txt.font.size = DocxPt(10.5)
        r_ex_txt.font.color.rgb = theme["text"]
        doc.add_paragraph().paragraph_format.space_after = DocxPt(12)

    # ── Sections ──────────────────────────────────────────────────────────────
    for sec in data.get("sections", []):
        # Section Heading
        h = doc.add_paragraph()
        h.paragraph_format.space_before = DocxPt(16)
        h.paragraph_format.space_after = DocxPt(6)
        r_h = h.add_run(sec.get("heading", "Section"))
        r_h.font.size = DocxPt(15)
        r_h.font.bold = True
        r_h.font.color.rgb = theme["heading"]

        # Section Intro Paragraph
        if sec.get("intro"):
            p_intro = doc.add_paragraph()
            p_intro.paragraph_format.space_after = DocxPt(8)
            p_intro.paragraph_format.line_spacing = 1.2
            r_intro = p_intro.add_run(sec.get("intro"))
            r_intro.font.size = DocxPt(11)
            r_intro.font.color.rgb = theme["text"]

        # Key Points (Bulleted with bold title)
        for pt in sec.get("key_points", []):
            bp = doc.add_paragraph(style="List Bullet")
            bp.paragraph_format.space_after = DocxPt(4)
            bp.paragraph_format.line_spacing = 1.15

            r_bt = bp.add_run(f"{pt.get('title')}: ")
            r_bt.font.bold = True
            r_bt.font.size = DocxPt(10.5)
            r_bt.font.color.rgb = theme["accent"]

            r_bd = bp.add_run(pt.get("description", ""))
            r_bd.font.size = DocxPt(10.5)
            r_bd.font.color.rgb = theme["text"]

        # Takeaway note
        if sec.get("takeaway"):
            p_tk = doc.add_paragraph()
            p_tk.paragraph_format.space_before = DocxPt(4)
            p_tk.paragraph_format.space_after = DocxPt(12)
            r_tk = p_tk.add_run(f"👉 {sec.get('takeaway')}")
            r_tk.font.size = DocxPt(10)
            r_tk.font.italic = True
            r_tk.font.color.rgb = theme["muted"]

    # ── Conclusion ────────────────────────────────────────────────────────────
    if data.get("conclusion"):
        hc = doc.add_paragraph()
        hc.paragraph_format.space_before = DocxPt(18)
        hc.paragraph_format.space_after = DocxPt(6)
        r_hc = hc.add_run("Strategic Summary & Recommendations")
        r_hc.font.size = DocxPt(15)
        r_hc.font.bold = True
        r_hc.font.color.rgb = theme["heading"]

        p_conc = doc.add_paragraph()
        p_conc.paragraph_format.line_spacing = 1.2
        r_conc = p_conc.add_run(data.get("conclusion"))
        r_conc.font.size = DocxPt(11)
        r_conc.font.color.rgb = theme["text"]

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(tmp.name)
    background_tasks.add_task(os.unlink, tmp.name)
    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=f"{req.topic}.docx",
        background=background_tasks,
    )


# ─── 3. PDF REPORT GENERATION (Consistent Multi-Page Report) ──────────────────
class DocuCraftPDF(FPDF):
    def __init__(self, theme_data, doc_title):
        super().__init__()
        self.theme = theme_data
        self.doc_title = clean_pdf_text(doc_title)
        self.is_dark = theme_data.get("is_dark", False)

    def header(self):
        # 1. Fill background of EVERY page automatically
        if self.is_dark:
            self.set_fill_color(*self.theme["bg"])
            self.rect(0, 0, 210, 297, "F")
        else:
            self.set_fill_color(*self.theme["bg"])
            self.rect(0, 0, 210, 297, "F")

        # 2. Running Header on Page 2+
        if self.page_no() > 1:
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(*self.theme["muted"])
            self.cell(0, 7, f"{self.doc_title}  |  DocuCraft Technical Report", 0, 0, "L")
            self.cell(0, 7, datetime.now().strftime("%B %Y"), 0, 1, "R")
            self.set_draw_color(*self.theme["border"])
            self.set_line_width(0.2)
            self.line(15, 14, 195, 14)
            self.ln(6)

    def footer(self):
        self.set_y(-14)
        self.set_font("Helvetica", "", 8)
        self.set_text_color(*self.theme["muted"])
        self.cell(0, 6, "DocuCraft AI Document Engine  -  Confidential", 0, 0, "L")
        self.cell(0, 6, f"Page {self.page_no()}", 0, 0, "R")


async def generate_pdf(req, background_tasks: BackgroundTasks):
    data = get_doc_ai_content(req.topic, req.slide_count)
    theme = PDF_THEMES.get(req.theme, PDF_THEMES["dark"])

    doc_title = clean_pdf_text(data.get("title", req.topic))
    pdf = DocuCraftPDF(theme, doc_title)
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(15, 15, 15)
    pdf.add_page()

    # ── Page 1 Header Banner ──────────────────────────────────────────────────
    # Category Pill
    pdf.set_font("Helvetica", "B", 9)
    pdf.set_text_color(*theme["accent"])
    pdf.cell(0, 6, "TECHNICAL WHITEPAPER & EXECUTIVE REPORT", ln=True)

    # Document Main Title
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(*theme["title"])
    pdf.multi_cell(180, 8.5, doc_title)
    pdf.ln(2)

    # Subtitle
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*theme["muted"])
    pdf.multi_cell(180, 5.5, clean_pdf_text(data.get("subtitle", "Comprehensive Technical Overview")))
    pdf.ln(3)

    # Decorative Divider Line
    pdf.set_draw_color(*theme["accent"])
    pdf.set_line_width(0.8)
    pdf.line(15, pdf.get_y(), 195, pdf.get_y())
    pdf.ln(5)

    # Executive Summary Card
    exec_summary = data.get("executive_summary", "")
    if exec_summary:
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_text_color(*theme["accent"])
        pdf.cell(0, 6, "EXECUTIVE SUMMARY", ln=True)

        pdf.set_font("Helvetica", "", 9.5)
        pdf.set_text_color(*theme["text"])
        pdf.multi_cell(180, 5, clean_pdf_text(exec_summary))
        pdf.ln(5)

    # ── Sections Rendering ────────────────────────────────────────────────────
    for sec in data.get("sections", []):
        # Section Heading
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*theme["heading"])
        pdf.cell(0, 7.5, clean_pdf_text(sec.get("heading", "Section")), ln=True)

        # Section Intro Paragraph
        if sec.get("intro"):
            pdf.set_font("Helvetica", "", 9.5)
            pdf.set_text_color(*theme["text"])
            pdf.multi_cell(180, 5, clean_pdf_text(sec.get("intro")))
            pdf.ln(2.5)

        # Key Points List
        for pt in sec.get("key_points", []):
            pdf.set_font("Helvetica", "B", 9.5)
            pdf.set_text_color(*theme["accent"])
            pdf.cell(5, 5, "-", 0, 0)
            pdf.cell(0, 5, clean_pdf_text(f"{pt.get('title')}: "), ln=True)

            pdf.set_font("Helvetica", "", 9.5)
            pdf.set_text_color(*theme["text"])
            pdf.set_x(20)
            pdf.multi_cell(175, 4.8, clean_pdf_text(pt.get("description", "")))
            pdf.ln(1.5)

        # Takeaway callout
        if sec.get("takeaway"):
            pdf.set_font("Helvetica", "I", 9)
            pdf.set_text_color(*theme["muted"])
            pdf.set_x(15)
            pdf.multi_cell(180, 4.8, clean_pdf_text(f">> {sec.get('takeaway')}"))
            pdf.ln(3)

        pdf.ln(2)

    # ── Conclusion ────────────────────────────────────────────────────────────
    if data.get("conclusion"):
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(*theme["heading"])
        pdf.cell(0, 6.5, "Summary & Strategic Takeaways", ln=True)

        pdf.set_font("Helvetica", "", 9.5)
        pdf.set_text_color(*theme["text"])
        pdf.multi_cell(180, 5, clean_pdf_text(data.get("conclusion")))

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf.output(tmp.name)
    background_tasks.add_task(os.unlink, tmp.name)
    return FileResponse(
        tmp.name,
        media_type="application/pdf",
        filename=f"{req.topic}.pdf",
        background=background_tasks,
    )